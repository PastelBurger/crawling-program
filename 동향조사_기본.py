import requests
import pandas as pd
import numpy as np
from bs4 import BeautifulSoup as bs
from pandas import DataFrame
from selenium import webdriver
import time
import openpyxl
import datetime
import matplotlib.pyplot as plt

url="C:/작업서류G/작업서류/1.기술특례상장/2.뷰노/"
file_name= "0.기초자료.xlsx"



df=pd.read_excel(url+file_name)

# ------------------------------------------------------------------------------------------------1. 연도별 출원 동향 그래프
df_1=df[["대표출원인",'출원일']]

s1 = []
time_format ="%Y.%m.%d"
for tag in df_1['출원일']:
      s1.append ( datetime.datetime.strptime ( tag , time_format ) )

df_1['출원일']= s1


df_1['출원연도']=df_1['출원일'].dt.strftime('%Y')

s1 = []
for tag in df_1['출원연도']:
      s1.append ( int(tag) )

df_1['출원연도']=s1

mn=min(s1)
mx=max(s1)

s1=[]
for tag in range(mn,mx+1) :
      ss=df_1['출원연도'][(df_1['출원연도'] == tag)].count ()
      s1.append(ss)

df_1_1=pd.DataFrame(zip(range(mn,mx+1),s1))
df_1_1.columns = ['출원연도', '개수']

df_1_11=pd.DataFrame.copy(df_1_1)
df_1_11.loc[len(df_1_11)]=['합계',sum(df_1_11['개수'])]

df_1_11.to_excel(url+'1.연도별출원동향(엑셀).xlsx')

# s1 = []
# for tag in df_1_1['출원연도']:
#       s1.append ( str(tag) )
#
# df_1_1['출원연도']=s1


# 한글 폰트 사용을 위해서 세팅
from matplotlib import font_manager, rc
font_path = "C:/Windows/Fonts/NGULIM.TTF"
font = font_manager.FontProperties(fname=font_path).get_name()
rc('font', family=font)

#-------------------------- 그래프그리기

plt.figure(figsize=(10, 5))
bar = plt.bar(df_1_1['출원연도'],df_1_1['개수'], color='#B69F5E')
plt.grid(True, axis='y',alpha=0.5, linestyle='--')
plt.xticks(df_1_1['출원연도'], rotation=45)
plt.gca().spines['right'].set_visible(False) #오른쪽 테두리 제거
plt.gca().spines['top'].set_visible(False) #위 테두리 제거
plt.gca().spines['left'].set_visible(False) #왼쪽 테두리 제거
# plt.gca().set_facecolor('#E6D4A1') #배경색
# plt.yticks(ticks= []) #y축 tick 제거

# 숫자 넣는 부분
for rect in bar:
    height = rect.get_height()
    plt.text(rect.get_x() + rect.get_width()/2.0, height,  height, ha='center', va='bottom', size = 10)

plt.savefig(url+'1.전체출원동향.jpg',dpi=300)

# ----------------------------------------------------------------------------------------1.1 연도별/국가별 출원 동향 그래프

df_1['국가코드']=df['국가코드']
nt= list(set(df_1['국가코드']))
num=len(df_1['국가코드'])

ss=[]
s1=[]
for tag in nt :
      ss=df_1['국가코드'][(df_1['국가코드'] == tag)].count ()
      s1.append(ss)

df_1_2= pd.DataFrame(nt)
df_1_2['개수']= s1
df_1_2.columns = ['국가', '개수'] # -------- 국가별 개수

ss=[]
s1=[]
s3=pd.DataFrame(range(mn,mx+1))
s3.columns = ['출원연도']
for tag2 in nt :
    ss = []
    s1 = []
    for tag in range(mn,mx+1) :
        ss=df_1['출원연도'][(df_1['출원연도'] == tag) & (df_1['국가코드'] == tag2) ].count ()
        s1.append(ss)
        s2 = pd.DataFrame ( zip ( range ( mn , mx + 1 ) , s1 ) )
        s2.columns = ['출원연도' , '개수']
        s3[tag2]=s2['개수']

df_1_3=s3 # ------- 국가별 연도별 출원 개수

df_1_2.to_excel(url+'2.국가별출원합계(엑셀).xlsx')
df_1_3.to_excel(url+'2.국가별출원동향(엑셀).xlsx')


#-------------------------- 그래프그리기

for tag in nt :
    plt.figure(figsize=(5, 5))
    plt.plot(df_1_1['출원연도'],df_1_1['개수'],'--', color='#B69F5E')
    plt.plot(df_1_3['출원연도'],df_1_3[tag], color='#B69F5E')
    plt.scatter( df_1_3['출원연도'] , df_1_3[tag] , color='#B69F5E',s=5)
    plt.fill_between(df_1_3['출원연도'],df_1_3[tag], 0, alpha=0.3, linewidth=0, edgecolor='grey', facecolor='grey', antialiased=True)
    plt.grid(True, axis='y',alpha=0.5, linestyle='--')
    # plt.xticks(df_1_1['출원연도'], rotation=45)
    # plt.gca().spines['right'].set_visible(False) #오른쪽 테두리 제거
    # plt.gca().spines['top'].set_visible(False) #위 테두리 제거
    # plt.gca().spines['left'].set_visible(False) #왼쪽 테두리 제거
    plt.savefig ( url + '2.국가별출원동향_' + tag + '.jpg' , dpi=300 )


# -----------------------------------------------------------------------------------------------------------2. 메인 IPC

# ipc=list(set(df['메인 IPC']))
# dt=len(df['메인 IPC'])
# s1=[]
# ss=[]
# for tag in ipc :
#       ss=df['메인 IPC'][(df['메인 IPC'] == tag)].count ()
#       s1.append(ss)

ipc=pd.DataFrame(df['메인 IPC'])
ipc.columns=['ipc']
# ipc['개수']= s1
# ipc.columns=['ipc','개수']

s1=[]
for tag in ipc['ipc'] :
    s1.append(tag[0:4])

ipc['대분류']= s1
df_2=ipc   # ----------IPC 분류 개수
# -----------------------------------
s1=[]
s1=list(set(df_2['대분류']))
s2=[]
ss=[]
for tag in s1 :
      ss=df_2['대분류'][(df_2['대분류'] == tag)].count ()
      s2.append(ss)

s3=pd.DataFrame(s1)
s3['개수']= s2
s3.columns=['대분류','개수']
s3=s3.sort_values('개수',ascending=False)
df_23=pd.DataFrame.copy(s3)

df_23.to_excel(url+'3.IPC대분류개수(엑셀).xlsx')

s3=s3.reset_index()
t=[]
t = sum(s3['개수'][range(4,len(s3['개수']))])
s3['개수'][4] =t

s4= s3.loc[0:4,:]
ss=[]
for tag in s4['개수'] :
    ss.append(round(tag/sum(s4['개수'])*100,2))

s4['퍼센트']=ss

df_2_1=s4
df_2_1['대분류'][4]='기타'
df_2_1=df_2_1.sort_values('개수',ascending=False)
df_2_1=df_2_1.reset_index()
df_2_1=df_2_1.loc[:, ['대분류','개수','퍼센트']]
df_2_1.columns=['대분류','개수','퍼센트']

# ------------------------------------
df24=df_2[df_2['대분류']==df_2_1['대분류'][0]]

ipc=list(set(df24['ipc']))
dt=len(df['메인 IPC'])
s1=[]
ss=[]
for tag in ipc :
      ss=df['메인 IPC'][(df['메인 IPC'] == tag)].count ()
      s1.append(ss)

ipc=pd.DataFrame(ipc)
ipc.columns=['ipc']
ipc['개수']= s1
ipc.columns=['ipc','개수']
s1=ipc
s1=s1.sort_values('개수',ascending=False)
s1=s1.reset_index()
ss=[]
for tag in s1['개수'] :
    ss.append(round(tag/sum(s1['개수'])*100,2))

s1['퍼센트']=ss
s2=s1.loc[0:3,:]
df_2_2=s2 # ----  iPC세부분류



#-------------------------- 그래프그리기
# 1. IPC 대분류 그리기
# 그래프 외곽선
wedgeprops = {
    'edgecolor': 'black',
    'linestyle': '-',
    'linewidth': 0.5,
    'width': 0.7
}
color=['#C58940','#D4A657','#E5BA73','#FAEAB1','#FAF8F1']

plt.figure(figsize=(5, 5))
plt.pie(df_2_1['개수'], labels=df_2_1['대분류'], startangle=90, explode=[ 0.1, 0, 0, 0,0], autopct='%.1f%%',counterclock=False,wedgeprops=wedgeprops,colors=color,shadow=True)
# plt.show()

plt.savefig(url+'3.IPC분류.jpg',dpi=300)

# 2. IPC 세부 분류 그리기
# df_2_2=sort_values('개수')
plt.figure(figsize=(5, 5))
bar = plt.bar(df_2_2['ipc'],df_2_2['개수'], color= color)
plt.grid(True, axis='y',alpha=0.5, linestyle='--')
# # plt.xticks(df_1_1['출원연도'], rotation=45)
plt.gca().spines['right'].set_visible(False) #오른쪽 테두리 제거
plt.gca().spines['top'].set_visible(False) #위 테두리 제거
# plt.gca().spines['left'].set_visible(False) #왼쪽 테두리 제거


# 숫자 넣는 부분
t=0
for rect in bar:
    height = rect.get_height()
    plt.text(rect.get_x() + rect.get_width()/2.0, height, str(df_2_2['퍼센트'][t])+'%', ha='center', va='bottom', size = 10)
    t=t+1

# plt.show()

plt.savefig(url+'3.1_IPC소분류.jpg',dpi=300)

# df_1.to_excel('./Test.xlsx', sheet_name='Sheet1')

# --------------------------------------------------------------------------------------------------------3. 공동출원 확인
# ------------------------------------------------------------------------------------------------------3.1 공동출원 출원동향
df3=df[df['대표출원인']!=df['출원인']]
s1 = []
time_format ="%Y.%m.%d"
for tag in df3['출원일']:
      s1.append ( datetime.datetime.strptime ( tag , time_format ) )

df3['출원일']= s1


df3['출원연도']=df3['출원일'].dt.strftime('%Y')

s1 = []
for tag in df3['출원연도']:
      s1.append ( int(tag) )

df3['출원연도']=s1
mn=min(s1)
mx=max(s1)
s1=[]
for tag in range(mn,mx+1) :
      ss=df3['출원연도'][(df3['출원연도'] == tag)].count ()
      s1.append(ss)

df31=pd.DataFrame(zip(range(mn,mx+1),s1))
df31.columns = ['출원연도', '개수']


# 한글 폰트 사용을 위해서 세팅
from matplotlib import font_manager, rc
font_path = "C:/Windows/Fonts/NGULIM.TTF"
font = font_manager.FontProperties(fname=font_path).get_name()
rc('font', family=font)

#-------------------------- 그래프그리기

plt.figure(figsize=(10, 5))
bar = plt.bar(df31['출원연도'],df31['개수'], color='#B69F5E')
plt.grid(True, axis='y',alpha=0.5, linestyle='--')
plt.xticks(df31['출원연도'], rotation=45)
plt.gca().spines['right'].set_visible(False) #오른쪽 테두리 제거
plt.gca().spines['top'].set_visible(False) #위 테두리 제거
plt.gca().spines['left'].set_visible(False) #왼쪽 테두리 제거
# plt.gca().set_facecolor('#E6D4A1') #배경색
# plt.yticks(ticks= []) #y축 tick 제거

# 숫자 넣는 부분
for rect in bar:
    height = rect.get_height()
    plt.text(rect.get_x() + rect.get_width()/2.0, height,  height, ha='center', va='bottom', size = 10)

plt.savefig(url+'3_1.공동출원동향.jpg',dpi=300)

# ---------------------------------------------------------------------------------------------------3.2 공동출원 출원인분석

alist=list(set(df3['출원인']))
s1=[]
ss=[]
for tag in alist:
    ss=df3['출원인'][(df3['출원인'] == tag)].count ()
    s1.append(ss)

df32=pd.DataFrame(alist)
df32['개수']=s1
df32.columns = ['공동출원인', '개수']
df32=df32.sort_values('개수',ascending=False)


if len(df32) == 1:
    df32.to_excel ( url + '3_2.공동출원분석('+str(len(df32))+'개사).xlsx' , sheet_name='Sheet1' )

if len(df32) == 2:
    df32.to_excel ( url + '3_2.공동출원분석('+str(len(df32))+'개사).xlsx' , sheet_name='Sheet1' )
    wedgeprops = {
        'edgecolor' : 'black' ,
        'linestyle' : '-' ,
        'linewidth' : 0.5 ,
        'width' : 0.7
    }
    color = ['#C58940' , '#E5BA73']
    plt.figure ( figsize=(5 , 5) )
    plt.pie ( df32['공동출원인'] , df32['개수'] , startangle=90 , explode=[0.1 , 0] , autopct='%.1f%%' ,
              counterclock=False , wedgeprops=wedgeprops , colors=color , shadow=True )
    plt.savefig (url + '3_2.공동출원분석('+str(len(df32))+'개사).jpg' , dpi=300 )

if len(df32) == 3:
    df32.to_excel ( url + '3_2.공동출원분석('+str(len(df32))+'개사).xlsx' , sheet_name='Sheet1' )
    wedgeprops = {
        'edgecolor' : 'black' ,
        'linestyle' : '-' ,
        'linewidth' : 0.5 ,
        'width' : 0.7
    }
    color = ['#C58940' , '#E5BA73','#FAEAB1']
    plt.figure ( figsize=(5 , 5) )
    plt.pie ( df32['공동출원인'] , df32['개수'] , startangle=90 , explode=[0.1 , 0,0] , autopct='%.1f%%' ,
              counterclock=False , wedgeprops=wedgeprops , colors=color , shadow=True )
    plt.savefig ( url + '3_2.공동출원분석('+str(len(df32))+'개사).jpg' , dpi=300 )

if len(df32) == 4:
    df32.to_excel ( url+'3_2.공동출원분석('+str(len(df32))+'개사).xlsx' , sheet_name='Sheet1' )
    wedgeprops = {
        'edgecolor' : 'black' ,
        'linestyle' : '-' ,
        'linewidth' : 0.5 ,
        'width' : 0.7
    }
    color = ['#C58940' , '#E5BA73' , '#FAEAB1','#FAF8F1']
    plt.figure ( figsize=(5 , 5) )
    plt.pie ( ddf32['공동출원인'] , df32['개수'] , startangle=90 , explode=[0.1 , 0 , 0] , autopct='%.1f%%' ,
              counterclock=False , wedgeprops=wedgeprops , colors=color , shadow=True )
    plt.savefig ( url + url + '3_2.공동출원분석('+str(len(df32))+'개사).jpg' , dpi=300 )

if len(df32) > 4:
    df32.to_excel ( url + '3_2.공동출원분석('+str(len(df32))+'개사).xlsx' , sheet_name='Sheet1' )
    ss = []
    for tag in df32['개수'] :
        ss.append ( round ( tag / sum ( df32['개수'] ) * 100 , 2 ) )
    df32['퍼센트'] = ss
    df32 = df32.reset_index ()
    df33 = df32.loc[0 :3 , :]
    color = ['#C58940' , '#E5BA73' , '#FAEAB1' , '#FAF8F1']
    # 그래프
    plt.figure ( figsize=(5 , 5) )
    bar = plt.bar ( df33['공동출원인'] , df33['개수'] , color=color )
    plt.grid ( True , axis='y' , alpha=0.5 , linestyle='--' )
    # plt.xticks( rotation=-45)
    plt.gca ().spines['right'].set_visible ( False )  # 오른쪽 테두리 제거
    plt.gca ().spines['top'].set_visible ( False )  # 위 테두리 제거
    # plt.gca().spines['left'].set_visible(False) #왼쪽 테두리 제거
    plt.gca ().axes.xaxis.set_visible ( False )  # x범위 없애기
    # 숫자 넣는 부분
    t = 0
    for rect in bar :
        height = rect.get_height ()
        plt.text ( rect.get_x () + rect.get_width () / 2.0 , height , str(height)+'('+str(df32['퍼센트'][t] ) + '%)' , ha='center', va='bottom' , size=10 )
        t = t + 1
    plt.savefig ( url + '3_2.공동출원분석('+str(len(df32))+'개사).jpg' , dpi=300 )

# ---------------------------------------------------------------------------------------------------3.3 공동출원 출원인_1,2위 추가분석
#-------------------------------------------------------------------1위
df331=df3[df3['출원인']==df33['공동출원인'][0]]

s1 = []
for tag in df331['출원연도']:
      s1.append ( int(tag) )

df331['출원연도']=s1
mn31=min(s1)
mx31=max(s1)
s1=[]
ss=[]
for tag in range(mn31,mx31+1) :
      ss=df331['출원연도'][(df331['출원연도'] == tag)].count ()
      s1.append(ss)

df3311=pd.DataFrame(zip(range(mn31,mx31+1),s1))
df3311.columns = ['출원연도', '개수']

#-------------------------- 그래프그리기

plt.figure(figsize=(7, 5))
bar = plt.bar(df3311['출원연도'],df3311['개수'], color='#B69F5E')
plt.grid(True, axis='y',alpha=0.5, linestyle='--')
plt.xticks(df3311['출원연도'], rotation=45)
plt.gca().spines['right'].set_visible(False) #오른쪽 테두리 제거
plt.gca().spines['top'].set_visible(False) #위 테두리 제거
plt.gca().spines['left'].set_visible(False) #왼쪽 테두리 제거
# plt.gca().set_facecolor('#E6D4A1') #배경색
# plt.yticks(ticks= []) #y축 tick 제거

# 숫자 넣는 부분
for rect in bar:
    height = rect.get_height()
    plt.text(rect.get_x() + rect.get_width()/2.0, height,  height, ha='center', va='bottom', size = 10)

plt.savefig(url+'3_3.공동출원1위 동향.jpg',dpi=300)

#-------------------------------------------------------------------2위

df332=df3[df3['출원인']==df33['공동출원인'][1]]

s1 = []
for tag in df332['출원연도']:
      s1.append ( int(tag) )

df332['출원연도']=s1
mn32=min(s1)
mx32=max(s1)
s1=[]
for tag in range(mn32,mx32+1) :
      ss=df332['출원연도'][(df332['출원연도'] == tag)].count ()
      s1.append(ss)

df3321=pd.DataFrame(zip(range(mn32,mx32+1),s1))
df3321.columns = ['출원연도', '개수']

#-------------------------- 그래프그리기

plt.figure(figsize=(7, 5))
bar = plt.bar(df3321['출원연도'],df3321['개수'], color='#B69F5E')
plt.grid(True, axis='y',alpha=0.5, linestyle='--')
plt.xticks(df3321['출원연도'], rotation=45)
plt.gca().spines['right'].set_visible(False) #오른쪽 테두리 제거
plt.gca().spines['top'].set_visible(False) #위 테두리 제거
plt.gca().spines['left'].set_visible(False) #왼쪽 테두리 제거
# plt.gca().set_facecolor('#E6D4A1') #배경색
# plt.yticks(ticks= []) #y축 tick 제거

# 숫자 넣는 부분
for rect in bar:
    height = rect.get_height()
    plt.text(rect.get_x() + rect.get_width()/2.0, height,  height, ha='center', va='bottom', size = 10)

plt.savefig(url+'3_3.공동출원2위 동향.jpg',dpi=300)

# -------------------------------------------------------------------------------------------------------4.주요 특허 선별(

df4=pd.DataFrame.copy(df)
df4=df4[~df['패밀리 문헌수'].isnull()]

mask=df4['패밀리 문헌수'].isin(['-'])
df4=df[~mask]

s1=[]
s2=[]
s1= df4.sort_values('패밀리 문헌수',ascending=False)
s1=s1.reset_index()

s2=s1.loc[0:20 ,]


df4=pd.DataFrame.copy(s2)

# ---------------------------------------------------------------------------------------------------4.1워드 클라우드(한글)

import nltk
from nltk.tokenize import word_tokenize
from collections import Counter
from wordcloud import WordCloud
from PIL import Image
from konlpy.tag import Okt

s1=df4['독립항'][0]
len(df4)

for tag in range(1,len(df4)) :
    s1= s1+df4['독립항'][tag]
    print(tag)

text = s1

# ----------------한글

okt = Okt()
line =[]
line = okt.pos(text)

n_adj =[] # 명사 또는 형용사인 단어만 n_adj에 넣어주기
for word, tag in line:
    if tag in ['Noun','Adjective']: # 명사 또는 형용사인 단어 및 2글자 이상인 단어 선택 시
        n_adj.append(word)

#제외할 단어 추가
stop_words = "상기 및 항 포함 를 장치 것 를 중 값 통해 부 있는 도 단계 내 제 대한 각각 로부터" #추가할 때 띄어쓰기로 추가해주기
stop_words = set(stop_words.split(' ')) # 불용어를 제외한 단어만 남기기
n_adj = [word for word in n_adj if not word in stop_words]

#가장 많이 나온 단어 100개 저장
counts = Counter(n_adj)
tags = counts.most_common(60)


# 이미지 추가(워드크라우드 모양 설정)
mask = Image.new("RGBA",(399,402), (255,255,255)) #(2555,2575)는 사진 크기, (255,255,255)는 색을의미
image = Image.open('C:/작업서류G/작업서류/9. 회사소개서/9.워드클라우드이미지/이미지1.jpg').convert("RGBA")
x,y = image.size
mask.paste(image,(0,0,x,y),image)
mask = np.array(mask)

#워드클라우드 특성
color='Wistia'
width=500
height=500
bcolor='white'
font='C:/Users/kaith/AppData/Local/Microsoft/Windows/Fonts/batang'
figsize=(5,5) # 워드크라워드 모양 바뀌면 바꿔 줘야함

wc = WordCloud(font_path= font, background_color=bcolor, mask=mask, colormap=color,width=width, height=height, scale=2.0, max_font_size=250)
gen = wc.generate_from_frequencies(dict(tags))


plt.figure(figsize=figsize)#사이즈 설정 및 출력
plt.imshow(gen,interpolation='bilinear')
plt.axis('off')# 차트로 나오지 않게
# plt.show()

plt.savefig(url+'4.1.한글워드클라우드.png',dpi=300)

# ---------------------------------------------------------------------------------------------------4.2워드 클라우드(영어)

df41=pd.DataFrame.copy(df)
df41=df41[df41['국가코드']=='US']
df41=df41[~df['패밀리 문헌수'].isnull()]
mask=df41['패밀리 문헌수'].isin(['-'])
df41=df41[~mask]

if len(df41) > 20 :
    s1=[]
    s2=[]
    s1= df41.sort_values('패밀리 문헌수',ascending=False)
    s1=s1.reset_index()
    s2=s1.loc[0:20 ,]
    df41 = pd.DataFrame.copy(s2)


df41 =df41.reset_index()
s1=df41['독립항'][0]
len(df41)


for tag in range(1,len(df41)) :
    s1= s1+df41['독립항'][tag]

text = s1

text1= word_tokenize(text)

stop_words1 = "청구 항 청구항1 . of a the and [ A at ] , ; an anda 's is by according or for based 1 step comprising about on corresponding in to least each among using that feature with method server" #추가할 때 띄어쓰기로 추가해주기
stop_words1 = set(stop_words1.split(' ')) # 불용어를 제외한 단어만 남기기
text2 = [word for word in text1 if not word in stop_words1]

#가장 많이 나온 단어 100개 저장
counts = Counter(text2)
tags = counts.most_common(60)

wc = WordCloud(font_path= font, background_color=bcolor, mask=mask, colormap=color,width=width, height=height, scale=2.0, max_font_size=250)
gen = wc.generate_from_frequencies(dict(tags))

plt.figure(figsize=figsize)#사이즈 설정 및 출력
plt.imshow(gen,interpolation='bilinear')
plt.axis('off')# 차트로 나오지 않게
# plt.show()

plt.savefig(url+'4.2.영어워드클라우드.png',dpi=300)

# df31.to_excel('./Test.xlsx', sheet_name='Sheet1')

# ------------------------------------------------------------------------------------------------------4.3.주요 특허 저장

from openpyxl import load_workbook
from openpyxl.styles import Alignment
from openpyxl.styles import PatternFill
from openpyxl.styles import colors
from openpyxl.styles import Color , Font, Border, Side


df42 = df4.loc[:,['출원번호','발명의 명칭','출원인','발명자','출원일','등록번호','법적상태','국가코드']]
df4.to_excel(url+'4.3.주요선별특허(엑셀).xlsx', sheet_name='Sheet1')
df42.to_excel(url+'4.3.주요선별특허(엑셀)_가공.xlsx', sheet_name='Sheet1')
# 엑셀 간격 조정
wb = load_workbook(url+'4.3.주요선별특허(엑셀)_가공.xlsx')
ws = wb['Sheet1']
ws.column_dimensions['A'].width=5
ws.column_dimensions['B'].width=13.25
ws.column_dimensions['C'].width=75.5
ws.column_dimensions['D'].width=20
ws.column_dimensions['E'].width=20
ws.column_dimensions['F'].width=10
ws.column_dimensions['G'].width=10
ws.column_dimensions['H'].width=10
ws.column_dimensions['I'].width=5

s1=[]
s1=['A','B','C','D','E','F','G','H','I']
ss=[]
ws['A1'].value='연번'
ws['A1'].font=Font(bold=True)

for tag in range(1,len(df42)+2) :
    for ss in range(0,9):
        a1 = ws[s1[ss] + str (tag)]
        a1.alignment = Alignment ( horizontal='left' , vertical='center' , wrap_text=True )
        a1.fill = PatternFill ( fill_type='solid' , fgColor=Color ( 'FAEAB1' ) )
        a1.border=Border ( left=Side ( border_style='thin' ,
                             color='000000' ) ,
                 right=Side ( border_style='thin' ,
                              color='000000' ) ,
                 top=Side ( border_style='thin' ,
                            color='000000' ) ,
                 bottom=Side ( border_style='thin' ,
                               color='000000' ) )


for ss in range(0,9): #1줄 변경
    a1 = ws[s1[ss] + str(1)]
    a1.alignment = Alignment ( horizontal='center' , vertical='center' , wrap_text=True )
    a1.fill = PatternFill ( fill_type='solid' , fgColor=Color ( 'E5BA73' ) )

for tag in range(1,len(df42)+2) : #45678 열 가운데 정렬
    for ss in range(5,9):
        a1 = ws[s1[ss] + str (tag)]
        a1.alignment = Alignment ( horizontal='center' , vertical='center' , wrap_text=True )
    for ss in range ( 0,2 ) :
        a1 = ws[s1[ss] + str ( tag )]
        a1.alignment = Alignment ( horizontal='center' , vertical='center' , wrap_text=True )




wb.save(url+'4.3.주요선별특허(엑셀)_가공.xlsx')


# ------------------------------------------------------------------------------------------------------4.4.국가별법적상태
df44=df.loc[: , ['국가코드','법적상태', '출원번호']]
df441=df44.groupby(['국가코드','법적상태'],as_index=False ).count()
df441.columns=['국가코드','법적상태','개수']

# df441.to_excel(url+'테스트.xlsx')
m1=list(set(df44['국가코드']))
m2=list(set(df44['법적상태']))
s1=[]
s2=[]
for tag in m1 :
    for tag2 in m2:
        s1.append(tag)
        s2.append(tag2)

df442=pd.DataFrame(s1)
df442['법적상태']=s2
df442.columns=['국가코드','법적상태']

ss=[]
s1=[]
for tag in range(0,len(df442)) :
      ss=df44['출원번호'][(df44['법적상태'] == df442['법적상태'][tag]) & (df44['국가코드'] == df442['국가코드'][tag])].count()
      s1.append(ss)

df442['개수']=s1





# --------------------------------------------그래프


def create_x(t, w, n, d):
    return [t*x + w*n for x in range(d)]

# t = 2 # Number of dataset
# w = 0.8 # We generally want bars to be 0.8
# n = 1 # A is first set of data
# d = 6 # Number of sets of bars
# x_values = [t*element + w*n for element in range(d)]
t=len(list(set(df442['법적상태']))) # Number of dataset
w=0.8
d=len(list(set(df442['국가코드']))) # Number of dataset
barn=[]
for tag in range(1,t+1): # 동적 변수 만들기
    barn.append(create_x(t,w,tag,d))

# 그래프 순서변경
s1=list(set(df442['법적상태']))
ss=['공개','심사중','소멸','등록']
s2=[]
for tag in range(0,len(ss)) :
    for tag1 in range(0,len(s1)) :
        if ss[tag] == s1[tag1] :
            s2.append(s1[tag1])

color=['#C58940','#E5BA73','#FAEAB1','#FAF8F1']
plt.grid ( True , axis='y' , alpha=0.5 , linestyle='--' )
ax = plt.subplot()
for tag in range(0,t) :
    s1= df442[df442['법적상태']== s2[tag]]
    ax.bar(barn[tag],s1['개수'], color=color[tag])
    # for rect in bar :
    #     height = rect.get_height ()
    #     plt.text ( rect.get_x () + rect.get_width () / 2.0 , height , height , ha='center' , va='bottom' ,
    #                size=10 )

# plt.xticks(df_1_1['출원연도'], rotation=45)
plt.gca ().spines['right'].set_visible ( False )  # 오른쪽 테두리 제거
plt.gca ().spines['top'].set_visible ( False )  # 위 테두리 제거
plt.gca ().spines['left'].set_visible ( False )  # 왼쪽 테두리 제거
# plt.gca().set_facecolor('#E6D4A1') #배경색
# plt.yticks(ticks= []) #y축 tick 제거


plt.savefig ( url + '테스트.jpg' , dpi=300 )






#
#
#
#
#
#
# s1=df442[df442['법적상태']== '공개']
# s2=df442[df442['법적상태']== '등록']
#
# bar=plt.bar ( s2['국가코드'] , s2['개수'] , color= '#FAEAB1')
# bar1=plt.bar ( s1['국가코드'] , s1['개수'] , color='#C58940',bottom=s2['개수'] )
# for rect in bar :
#     height = rect.get_height ()
#     plt.text ( rect.get_x () + rect.get_width () / 2.0 , height , height , ha='center' , va='bottom' , size=10 )
#
# for rect in bar1 :
#     height = rect.get_height ()
#     plt.text ( rect.get_x () + rect.get_width () / 2.0 , height , height , ha='center' , va='bottom' , size=10 )
#
#
# plt.savefig(url+'테스트.jpg',dpi=300)
#
#
# # ---------------------------그래프
#
# color=['#C58940','#D4A657','#E5BA73','#FAEAB1','#FAF8F1']
# t=0
#
#
#
# #
# # for tag in m1 :
# #     s1= df441[df441['법적상태']== tag]
# #     plt.figure ( figsize=(5 , 5) )
# #     ss=color[t]
# #     plt.bar ( s1['국가코드'] , s1['출원번호'] , color=ss )
# #     plt.grid ( True , axis='y' , alpha=0.5 , linestyle='--' )
# #     # plt.xticks(df_1_1['출원연도'], rotation=45)
# #     plt.gca ().spines['right'].set_visible ( False )  # 오른쪽 테두리 제거
# #     plt.gca ().spines['top'].set_visible ( False )  # 위 테두리 제거
# #     plt.gca ().spines['left'].set_visible ( False )  # 왼쪽 테두리 제거
# #     # plt.gca().set_facecolor('#E6D4A1') #배경색
# #     # plt.yticks(ticks= []) #y축 tick 제거
# #     # 숫자 넣는 부분
# #     for rect in bar :
# #         height = rect.get_height ()
# #         plt.text ( rect.get_x () + rect.get_width () / 2.0 , height , height , ha='center' , va='bottom' , size=10 )
#
#
# plt.show()
#
#
#
# df_1_11.to_excel(url+'1.연도별출원동향(엑셀).xlsx')



# ----------------------------------------------------------------------------------------------------------------ppt저장
#python-pptx의 가장 기본이 되는 Presentation 클래스 불러오기
# pip install python-pptx
from pptx import Presentation
from pptx.enum.text import PP_ALIGN   # 정렬 설정하기
from pptx.util import Pt      # Pt 폰트사이즈
from pptx.dml.color import RGBColor
from pptx.util import Cm

#Presentation 객체 생성
prs = Presentation(url+ 'Sample.pptx')
# -------------------------------------------------------------------------------------------슬라이드 1
# 슬라이드 지정하기
slide_num = 2
slide = prs.slides[slide_num]

# 슬라이드 내 shape 사전 만들기
shapes_list = slide.shapes
shape_index = {}
for i, shape in enumerate(shapes_list):
    shape_index[ shape.name ] = i

print(shape_index)   # {'Box_down': 0, 'Box_up': 1, 'name2': 2, 'name1': 3}

#텍스트 삽입 함수
def text_on_name1(shape, input_text,font_size = 18 , bold = True):
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alighnment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = input_text
    font = run.font
    font.color.rgb = RGBColor(182 , 159 , 94)
    font.size = Pt(font_size)
    font.bold = bold
    font.name = None

def text_on_name2(shape, input_text,font_size = 12 , bold = True):
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alighnment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = input_text
    font = run.font
    font.color.rgb = RGBColor(182 , 159 , 94)
    font.size = Pt(font_size)
    font.bold = bold
    font.name = None

# 텍스트 삽입

shape_name = 'name1'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name1(shape_select, '정량 분석 - 출원 동향')

shape_name = 'name2'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name2(shape_select, ' 연도별 출원 동향')

#이미지 추가
# 왼쪽, 높이 설정

left = Cm(4.58)
top = Cm(4.32)
height = Cm(9.51)

# 이미지 도형 추가
pic = slide.shapes.add_picture(url+"1.전체출원동향.jpg", left, top, height=height)

# -------------------------------------------------------------------------------------------슬라이드 2

# 슬라이드 지정하기
slide_num = 3
slide = prs.slides[slide_num]

# 슬라이드 내 shape 사전 만들기
shapes_list = slide.shapes
shape_index = {}
for i, shape in enumerate(shapes_list):
    shape_index[ shape.name ] = i

print(shape_index)   # {'Box_down': 0, 'Box_up': 1, 'name2': 2, 'name1': 3}

# 텍스트 삽입

shape_name = 'name1'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name1(shape_select, '정량 분석 - 출원 동향')

shape_name = 'name2'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name2(shape_select, ' 국가별 출원 동향')

#이미지 추가
# 이미지 도형 추가

# 1그림 좌표 6.98   4.73
# 2그림 좌표 15.59   4.73
# 3그림 좌표 6.98   8.89
# 4그림 좌표 15.59  8.89

if len(nt) == 4:
    t=0
    s1 = 6.17 , 12.86 , 6.17 , 12.86
    s1 = list ( s1 )
    s2 = 4.92 , 4.92 , 9.16 , 9.16
    s2 = list ( s2 )
    s3 = 'nation1' , 'nation2' , 'nation3' , 'nation4'
    s3 = list ( s3 )
    for tag in nt:
        shape_name = s3[t]
        shape_select = shapes_list[shape_index[shape_name]]
        t2='<'+tag+'>'
        text_on_name2 ( shape_select , t2 )
        left = Cm (s1[t])
        top = Cm ( s2[t])
        height = Cm ( 4.45 )
        pic = slide.shapes.add_picture ( url + '2.국가별출원동향_' + tag + '.jpg' , left , top , height=height )
        t=t+1


if len(nt) == 5:
    t=0
    s1 = 6.17 , 12.86 , 6.17 , 12.86 , 18.35
    s1 = list ( s1 )
    s2 = 4.92 , 4.92 , 9.16 , 9.16 , 6.6
    s2 = list ( s2 )
    s3 = 'nation1' , 'nation2' , 'nation3' , 'nation4', 'nation5'
    s3 = list ( s3 )
    for tag in nt:
        shape_name = s3[t]
        shape_select = shapes_list[shape_index[shape_name]]
        t2='<'+tag+'>'
        text_on_name2 ( shape_select , t2 )
        left = Cm (s1[t])
        top = Cm ( s2[t])
        height = Cm ( 4.45 )
        pic = slide.shapes.add_picture ( url + '2.국가별출원동향_' + tag + '.jpg' , left , top , height=height )
        t=t+1


# -------------------------------------------------------------------------------------------슬라이드 3
# 슬라이드 지정하기
slide_num = 4
slide = prs.slides[slide_num]

# 슬라이드 내 shape 사전 만들기
shapes_list = slide.shapes
shape_index = {}
for i, shape in enumerate(shapes_list):
    shape_index[ shape.name ] = i

print(shape_index)   # {'Box_down': 0, 'Box_up': 1, 'name2': 2, 'name1': 3}

# 텍스트 삽입
shape_name = 'name1'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name1(shape_select, '정량 분석 - IPC 분류')

shape_name = 'name2'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name2(shape_select, ' IPC 분류 분석')

#이미지1 추가

left = Cm(4.58)
top = Cm(4.32)
height = Cm(9.51)
pic = slide.shapes.add_picture(url+'3.IPC분류.jpg', left, top, height=height)

tx1='- '+'전체 IPC 분류 중'+ df_2_1['대분류'][0]+'이 '+str(df_2_1['퍼센트'][0])+'%로 1위로 많았고, '+ df_2_1['대분류'][1]+'이 '+str(df_2_1['퍼센트'][1])+'%로 2위로 많았음'
shape_name = 'text1'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name2(shape_select, tx1)

#이미지2 추가

left = Cm(14.58)
top = Cm(4.32)
height = Cm(9.51)
pic = slide.shapes.add_picture(url+'3.1_IPC소분류.jpg', left, top, height=height)

tx2='- '+ df_2_1['대분류'][0]+'중에서는'+ df_2_2['ipc'][0]+'이 '+str(df_2_2['퍼센트'][0])+'%로 가장 많았고, '+df_2_2['ipc'][1]+'이 '+str(df_2_2['퍼센트'][0])+'%로 두번째로 많았음'
shape_name = 'text2'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name2(shape_select, tx2)

# -------------------------------------------------------------------------------------------슬라이드 4

# 슬라이드 지정하기
slide_num = 5
slide = prs.slides[slide_num]

# 슬라이드 내 shape 사전 만들기
shapes_list = slide.shapes
shape_index = {}
for i, shape in enumerate(shapes_list):
    shape_index[ shape.name ] = i

print(shape_index)   # {'Box_down': 0, 'Box_up': 1, 'name2': 2, 'name1': 3}
# 텍스트 삽입

# 텍스트 삽입
shape_name = 'name1'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name1(shape_select, '정량 분석 - 공동 출원인 분석')

shape_name = 'name2'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name2(shape_select, ' 공동 출원인 분석 결과')

t=sum ( df32['개수'] )
tx1='- '+'전체 출원 중 공동출원인 건수는 '+ str(t)+ '건임. 그 중에서 출원인이"'+df33['공동출원인'][0]+'"인 건이 '+ str(df33['개수'][0])+'('+str(df33['퍼센트'][0])+'%)'+'로 가장 많음.'+'즉, "'+df33['공동출원인'][0]+'"과의 협업이 활발했던 것으로 예상됨.'

shape_name = 'text1'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name2(shape_select, tx1)

left = Cm(4.2)
top = Cm(6.33)
height = Cm(5.1)
pic = slide.shapes.add_picture(url+'3_1.공동출원동향.jpg', left, top, height=height)

left = Cm(14.11)
top = Cm(5.62)
height = Cm(5.83)
pic = slide.shapes.add_picture(url + '3_2.공동출원분석('+str(len(df32))+'개사).jpg', left, top, height=height)

# -------------------------------------------------------------------------------------------슬라이드 5

# 슬라이드 지정하기
slide_num = 6
slide = prs.slides[slide_num]

# 슬라이드 내 shape 사전 만들기
shapes_list = slide.shapes
shape_index = {}
for i, shape in enumerate(shapes_list):
    shape_index[ shape.name ] = i

print(shape_index)   # {'Box_down': 0, 'Box_up': 1, 'name2': 2, 'name1': 3}
# 텍스트 삽입

# 텍스트 삽입
shape_name = 'name1'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name1(shape_select, '정량 분석 - 공동 출원인 분석(1,2위 상세분석)')

shape_name = 'name2'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name2(shape_select, ' 1,2위 동향 분석 결과')

tx1='- '+'공동출원 1위인"'+df33['공동출원인'][0]+'" '+ str(mn31)+'년도부터 '+str(mn32)+'년도까지 협업한 것으로 판단되고, '+'공동출원 2위인"'+df33['공동출원인'][0]+'"과는 '+ str(mn32)+'년도부터 '+str(mㅌ32)+'년도까지 협업한 것으로 판단됨'

shape_name = 'text1'
shape_select = shapes_list[ shape_index[ shape_name ]]
text_on_name2(shape_select, tx1)

left = Cm(4.83)
top = Cm(5.77)
height = Cm(5.93)
pic = slide.shapes.add_picture(url+'3_3.공동출원1위 동향.jpg', left, top, height=height)

left = Cm(14.52)
top = Cm(5.77)
height = Cm(5.93)
pic = slide.shapes.add_picture(url+'3_3.공동출원1위 동향.jpg', left, top, height=height)

# -------------------------------------------------------------------------------------------슬라이드 6 - 워드클라우드

# 슬라이드 지정하기
slide_num = 7
slide = prs.slides[slide_num]

# 슬라이드 내 shape 사전 만들기
shapes_list = slide.shapes
shape_index = {}
for i, shape in enumerate(shapes_list):
    shape_index[ shape.name ] = i

print(shape_index)   # {'Box_down': 0, 'Box_up': 1, 'name2': 2, 'name1': 3}

#이미지 삽입

left = Cm(4.18)
top = Cm(2.82)
height = Cm(9.54)
pic = slide.shapes.add_picture(url+'4.1.한글워드클라우드.png', left, top, height=height)

left = Cm(13.83)
top = Cm(2.92)
height = Cm(9.5493)
pic = slide.shapes.add_picture(url+'4.2.영어워드클라우드.png', left, top, height=height)



# ------------------------------------------------------------------------------------------- Final_저장
prs.save(url+'test.pptx')